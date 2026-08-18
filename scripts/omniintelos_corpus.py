"""
OmniIntelOS document corpus generator.

Builds the enterprise document estate that sits alongside the KPI tables:
annual reports, board packs, an incident post-mortem, handbooks, policies,
technical whitepapers, meeting minutes, spreadsheets, slide decks and charts —
bilingual EN/FR, and grounded in the SAME 78-month model that produces the
KPIs, so a number quoted in a PDF matches the number in `kpi_metrics` for that
period. That agreement is the point: it lets the copilot be checked.

Imported by `scripts/seed_data.py`; not a separate script to run.

Optional dependencies degrade LOUDLY, never silently: if openpyxl /
python-pptx / matplotlib are missing, the formats they produce are reported as
skipped with the reason, rather than the corpus quietly shrinking.
"""
from __future__ import annotations

import json
import os
import textwrap
from pathlib import Path
from typing import Any, Dict, List, Optional

from omniintelos import (
    COMPANY, COMPANY_SHORT, HQ, PHASES, SERVICE_LINES, REGIONS, DEPARTMENTS,
    XOF_PER_USD, build_series, generate_kpis, health_index, months, phase_for,
)

try:
    from fpdf import FPDF
    _FPDF = True
except ImportError:
    _FPDF = False

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    _XLSX = True
except ImportError:
    _XLSX = False

try:
    from omniintelos_audio import build_audio, TTS_URL as _TTS_URL, TTS_TOKEN as _TTS_TOKEN
    _AUDIO = bool(_TTS_URL and _TTS_TOKEN)
except ImportError:
    _AUDIO = False  # TTS_ENDPOINT_URL/TTS_ENDPOINT_TOKEN not configured

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    _PPTX = True
except ImportError:
    _PPTX = False

try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    _MPL = True
except ImportError:
    _MPL = False


# ─────────────────────────────────────────────────────────────────────────────

def _l1(s: str) -> str:
    """fpdf2's core fonts are Latin-1. French accents are covered; typographic
    dashes and curly quotes are not, so fold them to ASCII equivalents rather
    than crashing mid-document."""
    return (s.replace("—", "-").replace("–", "-")
             .replace("’", "'").replace("‘", "'")
             .replace("“", '"').replace("”", '"')
             .replace("…", "...").replace(" ", " ")
             .replace("−", "-").replace("•", "-")
             .encode("latin-1", "replace").decode("latin-1"))


def usd(v: float) -> str:
    if abs(v) >= 1e6:
        return f"USD {v/1e6:,.2f}M"
    if abs(v) >= 1e3:
        return f"USD {v/1e3:,.1f}k"
    return f"USD {v:,.0f}"


class Doc(FPDF if _FPDF else object):
    """A4 report with running header/footer and a small typographic toolkit."""

    def __init__(self, title: str, subtitle: str = "", lang: str = "en"):
        super().__init__(orientation="P", unit="mm", format="A4")
        self.doc_title = _l1(title)
        self.doc_subtitle = _l1(subtitle)
        self.lang = lang
        self.set_auto_page_break(auto=True, margin=18)
        self.set_title(self.doc_title)
        self.set_author(COMPANY)

    def header(self):
        if self.page_no() == 1:
            return
        self.set_font("Helvetica", "", 7.5)
        self.set_text_color(120)
        self.cell(0, 6, _l1(f"{COMPANY_SHORT}  |  {self.doc_title}"), align="L")
        self.cell(0, 6, _l1(self.doc_subtitle), align="R", new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(210)
        self.line(self.l_margin, self.get_y() + 1, self.w - self.r_margin, self.get_y() + 1)
        self.ln(4)
        self.set_text_color(0)

    def footer(self):
        if self.page_no() == 1:
            return
        self.set_y(-14)
        self.set_font("Helvetica", "", 7.5)
        self.set_text_color(130)
        conf = "Document interne - diffusion restreinte" if self.lang == "fr" else "Internal - restricted distribution"
        self.cell(0, 5, _l1(conf), align="L")
        self.cell(0, 5, f"{self.page_no()}", align="R")
        self.set_text_color(0)

    # ── building blocks ──────────────────────────────────────────────────
    def _mc(self, h: float, txt: str, **kw):
        """multi_cell that always starts at the left margin and returns there.

        fpdf2's multi_cell defaults to new_x=XPos.LEFT — the left edge of the cell
        it just drew, not the page margin. Chained calls therefore walk rightwards
        until the usable width hits zero and fpdf raises 'Not enough horizontal
        space to render a single character'. Every text primitive below goes
        through here so that failure mode cannot come back."""
        self.set_x(self.l_margin)
        self.multi_cell(0, h, _l1(txt), new_x="LMARGIN", new_y="NEXT", **kw)

    def cover(self, meta: List[tuple]):
        self.add_page()
        self.ln(52)
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(37, 99, 235)
        self.cell(0, 7, _l1(COMPANY.upper()), new_x="LMARGIN", new_y="NEXT")
        self.set_text_color(0)
        self.set_font("Helvetica", "B", 25)
        self._mc(11, self.doc_title)
        self.ln(2)
        self.set_font("Helvetica", "", 13)
        self.set_text_color(90)
        self._mc(7, self.doc_subtitle)
        self.set_text_color(0)
        self.ln(14)
        self.set_draw_color(37, 99, 235)
        self.set_line_width(0.8)
        self.line(self.l_margin, self.get_y(), self.l_margin + 46, self.get_y())
        self.set_line_width(0.2)
        self.ln(12)
        self.set_font("Helvetica", "", 9.5)
        for k, v in meta:
            # multi_cell defaults to XPos.LEFT (the start of the cell just printed),
            # NOT the page margin — without an explicit LMARGIN reset the x position
            # drifts 46mm right on every row and runs off the page after four rows.
            self.set_x(self.l_margin)
            self.set_text_color(120)
            self.cell(46, 6, _l1(k))
            self.set_text_color(0)
            self.multi_cell(0, 6, _l1(str(v)), new_x="LMARGIN", new_y="NEXT")
        self.ln(20)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(140)
        disc = ("OmniIntelOS S.A. est une entreprise fictive. Ce document est genere "
                "pour la demonstration et l'evaluation d'IntelAI." if self.lang == "fr" else
                "OmniIntelOS S.A. is a fictional company. This document is generated for "
                "IntelAI demonstration and evaluation purposes.")
        self._mc(4.5, disc)
        self.set_text_color(0)

    def h1(self, t: str):
        if self.get_y() > self.h - 60:
            self.add_page()
        self.ln(4)
        self.set_font("Helvetica", "B", 15)
        self.set_text_color(17, 24, 39)
        self._mc(8, t)
        self.set_text_color(0)
        self.ln(1.5)

    def h2(self, t: str):
        if self.get_y() > self.h - 48:
            self.add_page()
        self.ln(2.5)
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(37, 99, 235)
        self._mc(6, t)
        self.set_text_color(0)
        self.ln(0.8)

    def para(self, t: str):
        self.set_font("Helvetica", "", 9.6)
        self._mc(5.1, t)
        self.ln(2.2)

    def bullets(self, items: List[str]):
        self.set_font("Helvetica", "", 9.6)
        for it in items:
            self.set_x(self.l_margin)
            self.cell(4.5, 5.1, "-")
            self.multi_cell(0, 5.1, _l1(it), new_x="LMARGIN", new_y="NEXT")
        self.ln(2)

    def kv_table(self, rows: List[tuple], widths=(78, 0)):
        self.set_font("Helvetica", "", 9.3)
        for k, v in rows:
            if self.get_y() > self.h - 26:
                self.add_page()
            self.set_fill_color(248, 250, 252)
            self.cell(widths[0], 6.4, _l1(f"  {k}"), border=0, fill=True)
            self.set_font("Helvetica", "B", 9.3)
            self.cell(0, 6.4, _l1(f"  {v}"), border=0, fill=True, new_x="LMARGIN", new_y="NEXT")
            self.set_font("Helvetica", "", 9.3)
            self.ln(0.7)
        self.ln(2)

    def table(self, headers: List[str], rows: List[List[str]], widths: List[float]):
        self.set_font("Helvetica", "B", 8.4)
        self.set_fill_color(37, 99, 235)
        self.set_text_color(255)
        for h, w in zip(headers, widths):
            self.cell(w, 6.6, _l1(f" {h}"), border=0, fill=True, align="L")
        self.ln()
        self.set_text_color(0)
        self.set_font("Helvetica", "", 8.4)
        for n, r in enumerate(rows):
            if self.get_y() > self.h - 24:
                self.add_page()
                self.set_font("Helvetica", "B", 8.4)
                self.set_fill_color(37, 99, 235)
                self.set_text_color(255)
                for h, w in zip(headers, widths):
                    self.cell(w, 6.6, _l1(f" {h}"), border=0, fill=True, align="L")
                self.ln()
                self.set_text_color(0)
                self.set_font("Helvetica", "", 8.4)
            self.set_fill_color(255 if n % 2 else 246, 255 if n % 2 else 249, 255 if n % 2 else 252)
            for c, w in zip(r, widths):
                self.cell(w, 5.9, _l1(f" {c}"), border=0, fill=True, align="L")
            self.ln()
        self.ln(3)

    def callout(self, title: str, body: str, rgb=(254, 243, 199)):
        if self.get_y() > self.h - 48:
            self.add_page()
        self.set_fill_color(*rgb)
        y0 = self.get_y()
        self.set_font("Helvetica", "B", 9.3)
        self._mc(5.6, f"  {title}", fill=True)
        self.set_font("Helvetica", "", 9.2)
        self._mc(5.0, textwrap.indent(body, "  "), fill=True)
        self.ln(3)


# ─────────────────────────────────────────────────────────────────────────────
# Data helpers
# ─────────────────────────────────────────────────────────────────────────────

# Healthy / risk bands, transcribed from IntelAI's published domain specification.
# Commentary is generated against THESE, so any statement of the form "below the
# X target" can be checked against both the stored KPI value and the spec.
#   metric -> (healthy_threshold, risk_threshold, higher_is_better, unit_suffix)
THRESHOLDS: Dict[str, tuple] = {
    "Gross Margin": (70.0, 25.0, True, "%"),
    "EBITDA Margin": (25.0, 5.0, True, "%"),
    "Net Profit Margin": (10.0, 0.0, True, "%"),
    "Cash Runway": (12.0, 4.0, True, " months"),
    "Debt to Equity": (1.5, 3.0, False, "x"),
    "Net Revenue Retention": (110.0, 90.0, True, "%"),
    "Monthly Churn Rate": (1.5, 5.0, False, "%"),
    "LTV to CAC Ratio": (3.0, 1.5, True, "x"),
    "CAC Payback Period": (12.0, 24.0, False, " months"),
    "Rule of 40": (40.0, 15.0, True, "%"),
    "Annual Employee Turnover": (10.0, 20.0, False, "%"),
    "Time to Hire": (35.0, 65.0, False, " days"),
    "Employee Net Promoter Score": (30.0, -10.0, True, ""),
    "Revenue per Employee": (300_000.0, 120_000.0, True, " USD"),
    "Offer Acceptance Rate": (85.0, 60.0, True, "%"),
    "Overall Equipment Effectiveness": (85.0, 60.0, True, "%"),
    "Defect Rate": (1.0, 4.5, False, "%"),
    "First Pass Yield": (92.0, 75.0, True, "%"),
    "Cycle Time Efficiency": (95.0, 80.0, True, "%"),
    "On-Time Delivery Rate": (95.0, 80.0, True, "%"),
    "Order Fulfillment Cycle Time": (48.0, 120.0, False, " hours"),
    "Inventory Turnover": (6.0, 2.5, True, "x"),
    "Supplier Defect Rate": (0.5, 3.0, False, "%"),
    "Carrying Cost of Inventory": (25.0, 35.0, False, "%"),
    "System Uptime": (99.95, 99.0, True, "%"),
    "Mean Time To Resolution": (0.5, 4.0, False, " hours"),
    "Change Failure Rate": (5.0, 15.0, False, "%"),
    "Critical Vulnerabilities": (0.0, 5.0, False, ""),
    "API P99 Latency": (250.0, 1500.0, False, " ms"),
    "Renewable Energy Ratio": (60.0, 20.0, True, "%"),
    "Board Diversity Ratio": (40.0, 15.0, True, "%"),
    "Audit Compliance Score": (98.0, 85.0, True, "%"),
    "Privacy Incident Count": (0.0, 0.5, False, ""),
}


def verdict(metric: str, value: float) -> str:
    """'on target' / 'below target' / 'in the risk band', judged against the
    published spec rather than against a feeling."""
    if metric not in THRESHOLDS:
        return "no published target"
    good, bad, higher, _ = THRESHOLDS[metric]
    if higher:
        if value >= good:
            return "on target"
        return "in the risk band" if value <= bad else "below target"
    if value <= good:
        return "on target"
    return "in the risk band" if value >= bad else "above target"


def fmt_metric(metric: str, value: float) -> str:
    suf = THRESHOLDS.get(metric, (0, 0, True, ""))[3]
    if suf == " USD":
        return usd(value)
    if abs(value) >= 10000:
        return f"{value:,.0f}{suf}"
    return f"{value:,.2f}{suf}".replace(".00", "")


def target_text(metric: str) -> str:
    if metric not in THRESHOLDS:
        return ""
    good, bad, higher, suf = THRESHOLDS[metric]
    arrow = ">=" if higher else "<="
    return f"target {arrow} {good:g}{suf}, risk at {'<=' if higher else '>='} {bad:g}{suf}"


def commentary(by, ms: List[str], metric: str, lang: str = "en") -> str:
    """One grounded analytical sentence about a metric over a window: level,
    direction, and standing against the published band. Reads differently for
    every metric and every period because it is computed, not templated prose."""
    vals = [by[m][metric] for m in ms if metric in by.get(m, {})]
    if not vals:
        return ""
    cur, first = vals[-1], vals[0]
    avg = sum(vals) / len(vals)
    v = verdict(metric, cur)
    delta = cur - first
    higher = THRESHOLDS.get(metric, (0, 0, True, ""))[2]
    improving = (delta > 0) == higher
    fr = lang == "fr"
    if fr:
        dirn = ("en amelioration" if improving else "en degradation") if abs(delta) > 1e-9 else "stable"
        vmap = {"on target": "conforme a la cible", "below target": "en deca de la cible",
                "above target": "au-dela de la cible", "in the risk band": "en zone de risque",
                "no published target": "sans cible publiee"}
        return (f"{metric} termine la periode a {fmt_metric(metric, cur)} "
                f"(moyenne {fmt_metric(metric, avg)}), {dirn} depuis {fmt_metric(metric, first)} - "
                f"{vmap[v]} ({target_text(metric)}).")
    dirn = ("improving" if improving else "deteriorating") if abs(delta) > 1e-9 else "flat"
    return (f"{metric} closed the period at {fmt_metric(metric, cur)} "
            f"(average {fmt_metric(metric, avg)}), {dirn} from {fmt_metric(metric, first)} - "
            f"{v} ({target_text(metric)}).")


DOMAIN_METRICS = {
    "Finance": ["Gross Margin", "EBITDA Margin", "Net Profit Margin", "Cash Runway", "Debt to Equity"],
    "Growth": ["Net Revenue Retention", "Monthly Churn Rate", "LTV to CAC Ratio",
               "CAC Payback Period", "Rule of 40"],
    "People": ["Annual Employee Turnover", "Time to Hire", "Employee Net Promoter Score",
               "Revenue per Employee", "Offer Acceptance Rate"],
    "Operations": ["Overall Equipment Effectiveness", "Defect Rate", "First Pass Yield",
                   "Cycle Time Efficiency"],
    "Logistics": ["On-Time Delivery Rate", "Order Fulfillment Cycle Time", "Inventory Turnover",
                  "Supplier Defect Rate", "Carrying Cost of Inventory"],
    "IT": ["System Uptime", "Mean Time To Resolution", "Change Failure Rate",
           "Critical Vulnerabilities", "API P99 Latency"],
    "ESG": ["Renewable Energy Ratio", "Board Diversity Ratio", "Audit Compliance Score",
            "Privacy Incident Count"],
}

# Revenue split used consistently across every document that reports segments.
REGION_MIX = {"Sahel": 0.34, "West Africa": 0.27, "North & East Africa": 0.12,
              "Europe": 0.19, "Americas": 0.08}
LINE_MIX = {"Data Science & Analytics": 0.24, "Computer Vision": 0.15,
            "NLP & Document Intelligence": 0.18, "IoT & Edge Telemetry": 0.09,
            "Blockchain & Provenance": 0.05, "Custom Software Engineering": 0.17,
            "Managed Data Centres": 0.12}


def _by_period(kpis) -> Dict[str, Dict[str, float]]:
    out: Dict[str, Dict[str, float]] = {}
    for r in kpis:
        out.setdefault(r["period"], {})[r["metric_name"]] = r["value"]
    return out


def _fy_months(year: int) -> List[str]:
    return [m for m in months() if m.startswith(str(year))]


def _agg(by: Dict[str, Dict[str, float]], ms: List[str], metric: str, how="sum") -> float:
    vals = [by[m][metric] for m in ms if m in by and metric in by[m]]
    if not vals:
        return 0.0
    if how == "sum":
        return sum(vals)
    if how == "avg":
        return sum(vals) / len(vals)
    if how == "last":
        return vals[-1]
    if how == "min":
        return min(vals)
    if how == "max":
        return max(vals)
    return 0.0


# ─────────────────────────────────────────────────────────────────────────────
# PDF builders
# ─────────────────────────────────────────────────────────────────────────────

def annual_report(by, year: int, lang: str, out: Path) -> Optional[Path]:
    ms = _fy_months(year)
    if not ms:
        return None
    fr = lang == "fr"
    rev = _agg(by, ms, "Revenue")
    cogs = _agg(by, ms, "COGS")
    opex = _agg(by, ms, "Operating Expenses")
    ebitda = _agg(by, ms, "EBITDA")
    ni = _agg(by, ms, "Net Income")
    arr_end = _agg(by, ms, "ARR", "last")
    hc_end = _agg(by, ms, "Headcount", "last")
    cust_end = _agg(by, ms, "Customers", "last")
    gm = (rev - cogs) / rev * 100 if rev else 0
    nrr = _agg(by, ms, "Net Revenue Retention", "avg")
    up = _agg(by, ms, "System Uptime", "avg")
    co2 = _agg(by, ms, "Total Carbon Footprint")
    ren = _agg(by, ms, "Renewable Energy Ratio", "last")
    turn = _agg(by, ms, "Annual Employee Turnover", "avg")
    enps = _agg(by, ms, "Employee Net Promoter Score", "avg")
    hscore, hlabel = health_index(ms[-1], by)
    ph = [p for p in PHASES if any(p["start"] <= m <= p["end"] for m in ms)]

    title = f"Rapport annuel {year}" if fr else f"Annual Report {year}"
    sub = ("Resultats consolides et revue des operations" if fr
           else "Consolidated results and review of operations")
    d = Doc(title, sub, lang)
    d.cover([
        ("Exercice" if fr else "Fiscal year", f"1 January - 31 December {year}"),
        ("Siege social" if fr else "Registered office", HQ),
        ("Forme juridique" if fr else "Legal form", "Societe Anonyme (OHADA)"),
        ("Monnaie de presentation" if fr else "Presentation currency", "USD (fonctionnelle: XOF)" if fr else "USD (functional: XOF)"),
        ("Date d'emission" if fr else "Issued", f"{year+1}-03-28"),
        ("Classification", "Interne" if fr else "Internal"),
    ])

    d.add_page()
    d.h1("Message du Directeur General" if fr else "Letter from the Chief Executive")
    for p in ph:
        d.para(p["narrative_fr"] if fr else p["narrative_en"])
    d.para(
        (f"Sur l'exercice, le chiffre d'affaires consolide s'etablit a {usd(rev)}, pour une marge brute "
         f"de {gm:.1f}% et un EBITDA de {usd(ebitda)}. L'ARR de sortie atteint {usd(arr_end)} aupres de "
         f"{cust_end:,.0f} clients actifs, avec un effectif de {hc_end:,.0f} collaborateurs. "
         f"L'indice de sante d'entreprise en fin d'exercice ressort a {hscore:.1f}/100 ({hlabel})."
         ) if fr else
        (f"Consolidated revenue for the year was {usd(rev)}, at a {gm:.1f}% gross margin and "
         f"{usd(ebitda)} of EBITDA. Exit ARR reached {usd(arr_end)} across {cust_end:,.0f} active "
         f"customers, with {hc_end:,.0f} employees. The enterprise health index closed the year at "
         f"{hscore:.1f}/100 ({hlabel}).")
    )

    d.h1("Chiffres cles" if fr else "Key figures")
    d.kv_table([
        ("Chiffre d'affaires" if fr else "Revenue", usd(rev)),
        ("Chiffre d'affaires (XOF)" if fr else "Revenue (XOF)", f"XOF {rev*XOF_PER_USD/1e9:,.2f} milliards" if fr else f"XOF {rev*XOF_PER_USD/1e9:,.2f} bn"),
        ("Marge brute" if fr else "Gross margin", f"{gm:.1f}%"),
        ("EBITDA", f"{usd(ebitda)} ({ebitda/rev*100 if rev else 0:.1f}%)"),
        ("Resultat net" if fr else "Net income", usd(ni)),
        ("ARR de sortie" if fr else "Exit ARR", usd(arr_end)),
        ("Retention nette du revenu" if fr else "Net revenue retention", f"{nrr:.1f}%"),
        ("Clients actifs" if fr else "Active customers", f"{cust_end:,.0f}"),
        ("Effectif" if fr else "Headcount", f"{hc_end:,.0f}"),
        ("Rotation du personnel" if fr else "Employee turnover", f"{turn:.1f}%"),
        ("eNPS", f"{enps:+.0f}"),
        ("Disponibilite des services" if fr else "Service uptime", f"{up:.3f}%"),
        ("Empreinte carbone" if fr else "Carbon footprint", f"{co2:,.0f} tCO2e"),
        ("Part d'energie renouvelable" if fr else "Renewable energy share", f"{ren:.1f}%"),
    ])

    d.h1("Compte de resultat consolide" if fr else "Consolidated statement of operations")
    d.table(
        ["Ligne" if fr else "Line item", "USD", "% CA" if fr else "% revenue"],
        [
            ["Chiffre d'affaires" if fr else "Revenue", f"{rev:,.0f}", "100.0%"],
            ["  dont abonnements" if fr else "  of which subscription", f"{_agg(by, ms, 'Subscription Revenue'):,.0f}", f"{_agg(by, ms, 'Subscription Revenue')/rev*100:.1f}%"],
            ["  dont services" if fr else "  of which services", f"{_agg(by, ms, 'Professional Services Revenue'):,.0f}", f"{_agg(by, ms, 'Professional Services Revenue')/rev*100:.1f}%"],
            ["  dont centre de donnees" if fr else "  of which data centre", f"{_agg(by, ms, 'Data Centre Revenue'):,.0f}", f"{_agg(by, ms, 'Data Centre Revenue')/rev*100:.1f}%"],
            ["Cout des ventes" if fr else "Cost of revenue", f"({cogs:,.0f})", f"{cogs/rev*100:.1f}%"],
            ["Marge brute" if fr else "Gross profit", f"{rev-cogs:,.0f}", f"{gm:.1f}%"],
            ["Charges operationnelles" if fr else "Operating expenses", f"({opex:,.0f})", f"{opex/rev*100:.1f}%"],
            ["EBITDA", f"{ebitda:,.0f}", f"{ebitda/rev*100:.1f}%"],
            ["Dotations aux amortissements" if fr else "Depreciation & amortisation", f"({_agg(by, ms, 'Depreciation & Amortisation'):,.0f})", ""],
            ["Charges financieres" if fr else "Interest expense", f"({_agg(by, ms, 'Interest Expense'):,.0f})", ""],
            ["Impots" if fr else "Taxes", f"({_agg(by, ms, 'Taxes'):,.0f})", ""],
            ["Resultat net" if fr else "Net income", f"{ni:,.0f}", f"{ni/rev*100:.1f}%"],
        ],
        [86, 48, 40],
    )
    d.para(
        ("Toutes les marges de ce tableau sont recalculables a partir des lignes qui le composent : "
         "la marge brute est (Chiffre d'affaires - Cout des ventes) / Chiffre d'affaires."
         ) if fr else
        ("Every margin in this table is re-derivable from the lines above it: gross margin is "
         "(Revenue - Cost of revenue) / Revenue."))

    d.h1("Performance mensuelle" if fr else "Monthly performance")
    d.table(
        ["Periode" if fr else "Period", "CA" if fr else "Revenue", "Marge brute" if fr else "Gross margin",
         "ARR", "Clients" if fr else "Customers", "Effectif" if fr else "Headcount", "Sante" if fr else "Health"],
        [[m, f"{by[m]['Revenue']:,.0f}", f"{by[m]['Gross Margin']:.1f}%",
          f"{by[m]['ARR']:,.0f}", f"{by[m]['Customers']:,.0f}",
          f"{by[m]['Headcount']:,.0f}", f"{health_index(m, by)[0]:.0f}"] for m in ms],
        [24, 30, 26, 30, 24, 24, 22],
    )

    # ── Segment reporting ────────────────────────────────────────────────
    d.h1("Information sectorielle" if fr else "Segment reporting")
    d.h2("Par region" if fr else "By region")
    d.para(
        ("La repartition geographique reflete l'implantation historique au Sahel et la montee en "
         "puissance des contrats europeens, portes par les obligations de residence des donnees et "
         "par les partenariats academiques."
         if fr else
         "The geographic split reflects the historic Sahel base and the growing European book, driven "
         "by data-residency obligations and academic partnerships."))
    d.table(["Region", "CA (USD)" if fr else "Revenue (USD)", "% CA" if fr else "% revenue",
             "CA (XOF)" if fr else "Revenue (XOF)"],
            [[rg, f"{rev*sh:,.0f}", f"{sh*100:.1f}%", f"{rev*sh*XOF_PER_USD/1e6:,.1f}M"]
             for rg, sh in REGION_MIX.items()],
            [56, 44, 30, 44])
    d.h2("Par ligne de service" if fr else "By service line")
    d.table(["Ligne de service" if fr else "Service line", "CA (USD)" if fr else "Revenue (USD)",
             "% CA" if fr else "% revenue"],
            [[ln, f"{rev*sh:,.0f}", f"{sh*100:.1f}%"] for ln, sh in LINE_MIX.items()],
            [86, 50, 38])

    # ── MD&A: computed commentary per domain, against the published bands ──
    d.h1("Analyse de la direction" if fr else "Management discussion and analysis")
    d.para(
        ("Chaque commentaire ci-dessous est calcule a partir des series mensuelles de l'exercice et "
         "confronte aux seuils publies dans la specification des domaines d'IntelAI. Un lecteur peut "
         "donc recalculer chaque affirmation a partir des donnees sous-jacentes."
         if fr else
         "Every comment below is computed from the year's monthly series and judged against the bands "
         "published in IntelAI's domain specification, so a reader can re-derive each statement from "
         "the underlying data."))
    for dom, mets in DOMAIN_METRICS.items():
        d.h2(dom)
        for met in mets:
            c = commentary(by, ms, met, lang)
            if c:
                d.para(c)

    d.h1("Revue par domaine" if fr else "Domain review")
    for dom, rows in [
        ("Croissance" if fr else "Growth", [
            ("ARR", usd(arr_end)), ("NRR", f"{nrr:.1f}%"),
            ("Taux d'attrition client mensuel" if fr else "Monthly logo churn", f"{_agg(by, ms, 'Monthly Churn Rate', 'avg'):.2f}%"),
            ("LTV:CAC", f"{_agg(by, ms, 'LTV to CAC Ratio', 'avg'):.1f}x"),
            ("Regle des 40" if fr else "Rule of 40", f"{_agg(by, ms, 'Rule of 40', 'avg'):.1f}%")]),
        ("Personnel" if fr else "People", [
            ("Effectif" if fr else "Headcount", f"{hc_end:,.0f}"),
            ("Rotation annuelle" if fr else "Annual turnover", f"{turn:.1f}%"),
            ("Delai de recrutement" if fr else "Time to hire", f"{_agg(by, ms, 'Time to Hire', 'avg'):.0f} " + ("jours" if fr else "days")),
            ("eNPS", f"{enps:+.0f}"),
            ("CA par employe" if fr else "Revenue per employee", usd(_agg(by, ms, 'Revenue per Employee', 'last')))]),
        ("Informatique et securite" if fr else "IT and security", [
            ("Disponibilite" if fr else "Uptime", f"{up:.3f}%"),
            ("MTTR", f"{_agg(by, ms, 'Mean Time To Resolution', 'avg'):.2f} h"),
            ("Taux d'echec des changements" if fr else "Change failure rate", f"{_agg(by, ms, 'Change Failure Rate', 'avg'):.1f}%"),
            ("Vulnerabilites critiques (max)" if fr else "Critical vulnerabilities (peak)", f"{_agg(by, ms, 'Critical Vulnerabilities', 'max'):.0f}"),
            ("Latence P99" if fr else "P99 latency", f"{_agg(by, ms, 'API P99 Latency', 'avg'):.0f} ms")]),
        ("Operations", [
            ("OEE", f"{_agg(by, ms, 'Overall Equipment Effectiveness', 'avg'):.1f}%"),
            ("Taux de defaut" if fr else "Defect rate", f"{_agg(by, ms, 'Defect Rate', 'avg'):.2f}%"),
            ("Rendement au premier passage" if fr else "First pass yield", f"{_agg(by, ms, 'First Pass Yield', 'avg'):.1f}%")]),
        ("Chaine logistique" if fr else "Supply chain", [
            ("Livraison a temps" if fr else "On-time delivery", f"{_agg(by, ms, 'On-Time Delivery Rate', 'avg'):.1f}%"),
            ("Rotation des stocks" if fr else "Inventory turnover", f"{_agg(by, ms, 'Inventory Turnover', 'avg'):.1f}x"),
            ("Taux de defaut fournisseur" if fr else "Supplier defect rate", f"{_agg(by, ms, 'Supplier Defect Rate', 'avg'):.2f}%")]),
        ("ESG", [
            ("Emissions totales" if fr else "Total emissions", f"{co2:,.0f} tCO2e"),
            ("Scope 1 / 2 / 3", f"{_agg(by, ms,'Scope 1 Emissions'):,.0f} / {_agg(by, ms,'Scope 2 Emissions'):,.0f} / {_agg(by, ms,'Scope 3 Emissions'):,.0f}"),
            ("Energie renouvelable" if fr else "Renewable energy", f"{ren:.1f}%"),
            ("Diversite du conseil" if fr else "Board diversity", f"{_agg(by, ms, 'Board Diversity Ratio', 'last'):.1f}%"),
            ("Score de conformite d'audit" if fr else "Audit compliance score", f"{_agg(by, ms, 'Audit Compliance Score', 'last'):.1f}%")]),
    ]:
        d.h2(dom)
        d.kv_table(rows)

    d.h1("Facteurs de risque" if fr else "Risk factors")
    d.bullets([
        ("Concentration geographique au Sahel : instabilite politique et fermetures de corridors "
         "peuvent retarder le materiel importe et la mobilite des equipes."
         if fr else
         "Sahel geographic concentration: political instability and corridor closures can delay "
         "imported hardware and team mobility."),
        ("Dependance energetique : la fiabilite du reseau conditionne les couts de groupes "
         "electrogenes et les emissions de Scope 1."
         if fr else
         "Energy dependence: grid reliability drives generator costs and Scope 1 emissions."),
        ("Risque de change : les revenus sont libelles en USD et EUR, les charges majoritairement "
         "en XOF (parite fixe EUR a 655,957)."
         if fr else
         "Currency risk: revenue is denominated in USD and EUR while costs are mostly XOF "
         "(fixed EUR peg at 655.957)."),
        ("Cybersecurite : l'incident de fevrier 2023 a demontre l'impact en cascade d'une "
         "intrusion sur la retention client et la tresorerie."
         if fr else
         "Cybersecurity: the February 2023 incident demonstrated how an intrusion cascades into "
         "customer retention and cash."),
        ("Penurie de competences : la concurrence regionale et internationale sur les profils IA "
         "maintient une pression sur la remuneration et l'attrition."
         if fr else
         "Talent scarcity: regional and international competition for AI skills keeps pressure on "
         "compensation and attrition."),
    ])

    # ── Notes to the accounts ────────────────────────────────────────────
    d.h1("Notes annexes" if fr else "Notes to the financial statements")
    notes = [
        ("Base de preparation" if fr else "Basis of preparation",
         ("Les comptes statutaires sont tenus en XOF selon le referentiel OHADA (SYSCOHADA revise). "
          "La presentation ci-dessus est convertie en USD au taux de planification interne, la parite "
          "EUR/XOF etant fixe a 655,957 par la BCEAO."
          if fr else
          "Statutory books are kept in XOF under the OHADA framework (revised SYSCOHADA). The "
          "presentation above is translated to USD at the internal planning rate; the EUR/XOF parity "
          "is fixed at 655.957 by the BCEAO.")),
        ("Reconnaissance du revenu" if fr else "Revenue recognition",
         ("Les abonnements sont reconnus lineairement sur la duree contractuelle. Les prestations de "
          "services sont reconnues a l'avancement. Les revenus d'hebergement sont reconnus sur la "
          "periode d'occupation."
          if fr else
          "Subscriptions are recognised rateably over the contract term. Professional services are "
          "recognised on a percentage-of-completion basis. Hosting revenue is recognised over the "
          "occupancy period.")),
        ("Immobilisations" if fr else "Property, plant and equipment",
         ("Le centre de donnees DC1 est amorti sur 15 ans pour le genie civil et 5 ans pour les "
          "equipements informatiques. Les accelerateurs graphiques sont amortis sur 4 ans."
          if fr else
          "The DC1 facility is depreciated over 15 years for civil works and 5 years for IT equipment. "
          "GPU accelerators are depreciated over 4 years.")),
        ("Impots" if fr else "Taxation",
         ("Impot sur les societes au Niger a 30%. Aucun actif d'impot differe n'est comptabilise sur "
          "les deficits reportables tant que leur recuperation n'est pas probable."
          if fr else
          "Niger corporate income tax at 30%. No deferred tax asset is recognised on carried-forward "
          "losses until recovery is probable.")),
        ("Engagements" if fr else "Commitments",
         ("Facilite d'equipement adossee au materiel du DC1, remboursable sur 60 mois. Engagements de "
          "capacite aupres de trois operateurs de transit."
          if fr else
          "Equipment facility secured on DC1 hardware, amortising over 60 months. Capacity commitments "
          "with three transit carriers.")),
        ("Evenements posterieurs" if fr else "Subsequent events",
         ("Aucun evenement significatif posterieur a la cloture n'est a signaler."
          if fr else
          "No material events after the reporting date are reported.")),
    ]
    for i, (t, b) in enumerate(notes, 1):
        d.h2(f"{i}. {t}")
        d.para(b)

    d.h1("Gouvernance" if fr else "Governance")
    d.kv_table([
        ("Conseil d'administration" if fr else "Board of directors", "7 " + ("membres, dont 3 independants" if fr else "members, 3 independent")),
        ("Comite d'audit" if fr else "Audit committee", "3 " + ("membres, president independant" if fr else "members, independent chair")),
        ("Diversite du conseil" if fr else "Board diversity", f"{_agg(by, ms, 'Board Diversity Ratio', 'last'):.1f}%"),
        ("Commissaire aux comptes" if fr else "Statutory auditor", "Cabinet inscrit ONECCA Niger" if fr else "ONECCA Niger registered firm"),
        ("Score de conformite d'audit" if fr else "Audit compliance score", f"{_agg(by, ms, 'Audit Compliance Score', 'last'):.1f}%"),
        ("Incidents de confidentialite" if fr else "Privacy incidents", f"{_agg(by, ms, 'Privacy Incident Count'):.0f}"),
    ])

    # ── Full data appendix: every metric, every month of the year ────────
    d.h1("Annexe - donnees mensuelles detaillees" if fr else "Appendix - detailed monthly data")
    d.para(
        ("Cette annexe reproduit l'integralite des indicateurs mensuels de l'exercice, domaine par "
         "domaine. Elle existe pour que toute affirmation du present rapport puisse etre verifiee "
         "ligne par ligne contre les donnees sources."
         if fr else
         "This appendix reproduces every monthly indicator for the year, domain by domain. It exists so "
         "that any statement in this report can be checked line by line against the source data."))
    kpis_all = generate_kpis()
    per_dom: Dict[str, List[str]] = {}
    for r in kpis_all:
        per_dom.setdefault(r["category"], [])
        if r["metric_name"] not in per_dom[r["category"]]:
            per_dom[r["category"]].append(r["metric_name"])
    for dom in ["Finance", "Growth", "People", "Operations", "Logistics", "IT", "ESG"]:
        mets = per_dom.get(dom, [])
        if not mets:
            continue
        d.h2(dom)
        # Chunk metrics so each table stays inside the page width.
        for chunk_start in range(0, len(mets), 5):
            chunk = mets[chunk_start:chunk_start + 5]
            d.table(["Period"] + [m[:22] for m in chunk],
                    [[m] + [f"{by[m].get(x, 0):,.2f}" for x in chunk] for m in ms],
                    [22] + [33] * len(chunk))

    d.h1("Perspectives" if fr else "Outlook")
    nxt = [p for p in PHASES if p["start"] > f"{year}-12"]
    if nxt:
        d.para(nxt[0]["narrative_fr"] if fr else nxt[0]["narrative_en"])
    else:
        d.para("La direction maintient ses priorites : croissance disciplinee, marge et certification ESG."
               if fr else
               "Management maintains its priorities: disciplined growth, margin, and ESG certification.")

    p = out / f"omniintelos_annual_report_{year}_{lang}.pdf"
    d.output(str(p))
    return p


def incident_postmortem(by, out: Path) -> Path:
    d = Doc("Security Incident Post-Mortem", "INC-2023-0214 - Staging estate intrusion", "en")
    d.cover([
        ("Incident ID", "INC-2023-0214"),
        ("Severity", "SEV-1 (Critical)"),
        ("Detected", "2023-02-14 02:41 UTC"),
        ("Contained", "2023-02-16 19:05 UTC"),
        ("Closed", "2023-05-31"),
        ("Author", "Security & Compliance, OmniIntelOS S.A."),
        ("Classification", "Internal - restricted"),
    ])
    d.add_page()
    d.h1("1. Executive summary")
    d.para(
        "Between 13 and 16 February 2023 an external actor obtained valid credentials to a staging "
        "environment that held a replicated subset of customer telemetry. The intrusion was detected by "
        "anomalous egress volume, contained within 65 hours, and did not reach production customer "
        "databases. Three privacy incidents were formally notified to affected customers and to the "
        "supervisory authority.")
    d.callout("Measured impact",
              f"February 2023 uptime fell to {by['2023-02']['System Uptime']:.2f}% against a 99.95% target. "
              f"Critical vulnerabilities peaked at {by['2023-02']['Critical Vulnerabilities']:.0f}. "
              f"Monthly logo churn rose to {by['2023-04']['Monthly Churn Rate']:.2f}% by April as affected "
              f"customers exercised termination rights, and net revenue retention bottomed at "
              f"{min(by[m]['Net Revenue Retention'] for m in ['2023-02','2023-03','2023-04','2023-05']):.1f}%.",
              (254, 226, 226))

    d.h1("2. Timeline")
    d.table(["Timestamp (UTC)", "Event"], [
        ["2023-02-13 21:14", "Credential-stuffing attempts begin against the staging SSO endpoint."],
        ["2023-02-13 23:02", "Valid session established using a re-used contractor credential; MFA not enforced on staging."],
        ["2023-02-14 02:41", "Egress anomaly detection fires: 4.2 GB outbound from a staging subnet."],
        ["2023-02-14 03:10", "On-call SRE escalates to SEV-1; incident channel opened."],
        ["2023-02-14 05:55", "Staging estate isolated from the transit VPC. Customer-facing degradation begins."],
        ["2023-02-14 11:20", "Forensic imaging starts; external DFIR retained."],
        ["2023-02-15 08:00", "Scope confirmed: staging telemetry replica only; production RDS untouched."],
        ["2023-02-15 16:30", "Global credential rotation; MFA enforced on all non-production estates."],
        ["2023-02-16 19:05", "Containment declared. Phased service restoration begins."],
        ["2023-02-20", "Customer notifications dispatched (EN/FR); regulator notified within 72h."],
        ["2023-03-06", "Independent penetration test commissioned."],
        ["2023-05-31", "All 14 remediation actions closed; incident formally closed."],
    ], [34, 144])

    d.h1("3. Root cause")
    d.para(
        "Three independent control failures had to coincide. None alone would have produced the incident, "
        "which is why none was caught by the controls that existed at the time.")
    d.bullets([
        "Non-production estates were exempt from the MFA policy that covered production, on the assumption "
        "that staging held no customer data. That assumption had silently stopped being true when a "
        "telemetry replication job was pointed at staging in Q3 2022.",
        "A contractor credential issued in 2021 was never revoked at contract end; the joiner-mover-leaver "
        "process covered employees but not contractors.",
        "Egress monitoring alerted on volume but not on destination reputation, so the first 3.5 hours of "
        "exfiltration to a known-bad ASN did not raise an alert.",
    ])

    d.h1("4. Cascade into other domains")
    d.para(
        "The incident is the clearest example in the company's history of single-domain failure propagating "
        "across the business, and the monthly KPI series records each step.")
    d.table(["Month", "Uptime %", "SLA %", "Logo churn %", "NRR %", "Revenue (USD)", "Health"],
            [[m, f"{by[m]['System Uptime']:.2f}", f"{by[m]['SLA Compliance']:.1f}",
              f"{by[m]['Monthly Churn Rate']:.2f}", f"{by[m]['Net Revenue Retention']:.1f}",
              f"{by[m]['Revenue']:,.0f}", f"{health_index(m, by)[0]:.0f}"]
             for m in ["2023-01", "2023-02", "2023-03", "2023-04", "2023-05", "2023-06", "2023-07"]],
            [22, 26, 24, 28, 24, 34, 20])

    d.h1("5. Remediation actions")
    d.table(["#", "Action", "Owner", "Status"], [
        ["1", "Enforce MFA on every estate including non-production", "Security", "Closed 2023-02-15"],
        ["2", "Extend joiner-mover-leaver to contractors and service accounts", "People + Security", "Closed 2023-03-02"],
        ["3", "Destination-reputation egress filtering", "Cloud Ops", "Closed 2023-03-17"],
        ["4", "Prohibit customer data in non-production; automated scanner", "Data Platform", "Closed 2023-04-11"],
        ["5", "Zero-trust segmentation between estates", "Cloud Ops", "Closed 2023-05-08"],
        ["6", "Independent penetration test, annual cadence", "Security", "Closed 2023-05-22"],
        ["7", "ISO 27001 programme initiation", "Compliance", "Closed 2023-05-31"],
    ], [10, 96, 42, 40])

    d.h1("6. What we would not change")
    d.para(
        "Isolating the staging estate at 05:55 on 14 February caused most of the measured customer-facing "
        "degradation that month. It was the correct call and would be made again: the alternative was "
        "leaving a confirmed active intrusion connected to the transit network while forensics ran.")
    p = out / "omniintelos_incident_postmortem_INC-2023-0214_en.pdf"
    d.output(str(p))
    return p


def employee_handbook(out: Path, lang: str) -> Path:
    fr = lang == "fr"
    d = Doc("Manuel de l'employe" if fr else "Employee Handbook",
            "Edition 2026.1 - Toutes entites" if fr else "Edition 2026.1 - All entities", lang)
    d.cover([
        ("Version", "2026.1"),
        ("Entree en vigueur" if fr else "Effective", "2026-01-15"),
        ("Perimetre" if fr else "Scope", "Niger, Mali, Burkina Faso, Tchad, Senegal, France, USA"),
        ("Proprietaire" if fr else "Owner", "People & Culture"),
        ("Revision" if fr else "Review cycle", "Annuelle" if fr else "Annual"),
    ])
    secs = [
        ("Bienvenue chez OmniIntelOS" if fr else "Welcome to OmniIntelOS",
         ["Notre mission est de rendre l'intelligence artificielle appliquee accessible aux organisations "
          "africaines, avec des equipes basees en Afrique et des standards d'ingenierie internationaux."
          if fr else
          "Our mission is to make applied artificial intelligence accessible to African organisations, "
          "with teams based in Africa and international engineering standards.",
          "Ce manuel s'applique a l'ensemble des entites du groupe. En cas de divergence avec le droit "
          "local du travail, le droit local prevaut."
          if fr else
          "This handbook applies across all group entities. Where it conflicts with local labour law, "
          "local law prevails."]),
        ("Langues de travail" if fr else "Working languages",
         ["Le francais et l'anglais sont les deux langues officielles. Toute documentation destinee a "
          "l'ensemble du groupe est publiee dans les deux langues. Les reunions regionales du Sahel se "
          "tiennent en francais ; les revues d'architecture et les comites produits se tiennent en anglais."
          if fr else
          "French and English are both official. Any documentation intended for the whole group is "
          "published in both. Sahel regional meetings are held in French; architecture reviews and "
          "product councils are held in English."]),
        ("Organisation du travail" if fr else "Ways of working",
         ["Modele hybride : trois jours par semaine sur site pour les equipes de livraison, flexibilite "
          "totale pour les fonctions transverses." if fr else
          "Hybrid model: three days per week on site for delivery teams, full flexibility for shared functions.",
          "Les astreintes concernent les equipes Cloud & Data Centre Operations et Security ; elles sont "
          "compensees selon la grille en annexe B." if fr else
          "On-call applies to Cloud & Data Centre Operations and Security teams and is compensated per the "
          "schedule in Annex B.",
          "Semaine standard de 40 heures. Les heures supplementaires sont approuvees a l'avance par le "
          "responsable hierarchique." if fr else
          "Standard 40-hour week. Overtime is approved in advance by the line manager."]),
        ("Remuneration et avantages" if fr else "Compensation and benefits",
         ["Les grilles salariales sont revues annuellement et referencees sur des enquetes de marche "
          "regionales, avec un ecart cible inferieur a 5% entre profils equivalents." if fr else
          "Salary bands are reviewed annually against regional market surveys, targeting under 5% spread "
          "between equivalent profiles.",
          "Couverture sante pour l'employe et les ayants droit directs, prise en charge a 80%." if fr else
          "Health cover for the employee and direct dependants, 80% employer-funded.",
          "Plan d'interessement en actions ouvert a tous les employes apres 12 mois d'anciennete." if fr else
          "Equity participation plan open to all employees after 12 months of service."]),
        ("Securite de l'information" if fr else "Information security",
         ["L'authentification multifacteur est obligatoire sur TOUS les environnements, production comme "
          "hors production. Cette regle est issue directement de l'incident INC-2023-0214." if fr else
          "Multi-factor authentication is mandatory on ALL environments, production and non-production "
          "alike. This rule comes directly from incident INC-2023-0214.",
          "Aucune donnee client ne doit etre copiee vers un environnement hors production. Un scanner "
          "automatise verifie cette regle quotidiennement." if fr else
          "No customer data may be copied into a non-production environment. An automated scanner verifies "
          "this daily.",
          "Tout acces est revoque au plus tard le dernier jour travaille, employes et prestataires compris."
          if fr else
          "All access is revoked no later than the last working day, for employees and contractors alike."]),
        ("Conduite et ethique" if fr else "Conduct and ethics",
         ["Tolerance zero pour le harcelement, la discrimination et la corruption. Un canal de signalement "
          "anonyme est disponible et gere par un tiers independant." if fr else
          "Zero tolerance for harassment, discrimination and bribery. An anonymous reporting channel is "
          "available and administered by an independent third party.",
          "Les cadeaux et invitations d'une valeur superieure a 50 000 XOF doivent etre declares." if fr else
          "Gifts and hospitality above XOF 50,000 must be declared."]),
        ("Conges" if fr else "Leave",
         ["30 jours ouvrables de conges payes par an, plus les jours feries locaux de l'entite de "
          "rattachement." if fr else
          "30 working days of paid leave per year, plus the local public holidays of the employing entity.",
          "Conge maternite : 14 semaines minimum ; conge du second parent : 4 semaines, entierement "
          "remunerees." if fr else
          "Maternity leave: 14 weeks minimum; second-parent leave: 4 weeks, fully paid."]),
        ("Developpement professionnel" if fr else "Professional development",
         ["Budget annuel de formation par employe, avec priorite aux certifications cloud, securite et "
          "science des donnees." if fr else
          "Annual training budget per employee, prioritising cloud, security and data science certifications.",
          "Partenariats universitaires (Universite Abdou Moumouni, UCAD) pour l'encadrement de stages et "
          "de theses appliquees." if fr else
          "University partnerships (Universite Abdou Moumouni, UCAD) hosting internships and applied theses."]),
        ("Recrutement et integration" if fr else "Recruitment and onboarding",
         ["Tout poste ouvert fait l'objet d'une grille d'evaluation structuree, identique pour tous les "
          "candidats, afin de limiter le biais d'entretien." if fr else
          "Every open role uses a structured scorecard, identical across candidates, to limit interview bias.",
          "L'integration dure 90 jours et comporte un parrain designe, un plan de montee en competence "
          "et trois points formels a 30, 60 et 90 jours." if fr else
          "Onboarding runs 90 days with a named buddy, a ramp plan, and formal checkpoints at 30, 60 and 90 days.",
          "La capacite d'integration est plafonnee : le nombre de recrutements simultanes ne peut depasser "
          "la capacite d'encadrement disponible. Cette regle decoule directement de la crise de croissance "
          "de 2022." if fr else
          "Onboarding capacity is capped: simultaneous hires may not exceed available mentoring capacity. "
          "This rule comes directly from the 2022 growing-pains period."]),
        ("Evaluation de la performance" if fr else "Performance management",
         ["Cycle semestriel avec objectifs ecrits, auto-evaluation, evaluation par le responsable et "
          "calibration collective." if fr else
          "Half-yearly cycle with written objectives, self-assessment, manager assessment and group calibration.",
          "Les augmentations et promotions sont decidees en calibration, jamais unilateralement par un "
          "seul responsable." if fr else
          "Increases and promotions are decided in calibration, never unilaterally by a single manager."]),
        ("Mobilite et voyages" if fr else "Mobility and travel",
         ["Les deplacements entre entites du Sahel requierent une validation securite prealable, revue "
          "mensuellement selon la situation regionale." if fr else
          "Travel between Sahel entities requires prior security clearance, reviewed monthly against the "
          "regional situation.",
          "Les voyages internationaux sont regroupes lorsque possible afin de limiter les emissions de "
          "Scope 3." if fr else
          "International travel is batched where possible to limit Scope 3 emissions."]),
        ("Materiel et acces" if fr else "Equipment and access",
         ["Chaque employe recoit un poste chiffre au repos, gere par la flotte, avec verrouillage "
          "automatique et effacement a distance." if fr else
          "Each employee receives a fleet-managed laptop with full-disk encryption, automatic lock and "
          "remote wipe.",
          "L'utilisation d'appareils personnels pour acceder aux donnees classees Restreint est interdite."
          if fr else
          "Use of personal devices to access data classified Restricted is prohibited."]),
        ("Sante, securite et bien-etre" if fr else "Health, safety and wellbeing",
         ["Plan de continuite pour les sites exposes, exercices d'evacuation semestriels et dispositif "
          "d'assistance en cas d'incident regional." if fr else
          "Continuity plan for exposed sites, half-yearly evacuation drills, and an assistance scheme for "
          "regional incidents.",
          "Soutien psychologique confidentiel pris en charge, accessible a l'employe et a sa famille "
          "proche." if fr else
          "Employer-funded confidential counselling, available to the employee and immediate family."]),
        ("Depart de l'entreprise" if fr else "Leaving the company",
         ["Preavis conforme au droit local. Entretien de depart systematique dont les themes recurrents "
          "sont rapportes trimestriellement au comite de direction." if fr else
          "Notice per local law. A structured exit interview is held in every case, with recurring themes "
          "reported quarterly to the executive committee.",
          "Revocation de tous les acces au plus tard le dernier jour travaille, sans exception." if fr else
          "All access revoked no later than the last working day, without exception."]),
    ]
    d.add_page()
    for i, (title, paras) in enumerate(secs, 1):
        d.h1(f"{i}. {title}")
        for pp in paras:
            d.para(pp)
    d.h1(("Annexe A - Departements" if fr else "Annex A - Departments"))
    d.bullets(DEPARTMENTS)
    d.h1(("Annexe B - Grille d'astreinte" if fr else "Annex B - On-call schedule"))
    d.table(["Tier", "Response target", "Compensation"], [
        ["SEV-1", "15 minutes, 24/7", "Premium rate x2.0"],
        ["SEV-2", "1 hour, business + extended", "Premium rate x1.5"],
        ["SEV-3", "Next business day", "Standard rate"],
    ], [40, 76, 60])
    p = out / f"omniintelos_employee_handbook_2026_{lang}.pdf"
    d.output(str(p))
    return p


def datacentre_whitepaper(by, out: Path) -> Path:
    d = Doc("Niamey DC1 Technical Whitepaper", "Sovereign AI infrastructure for the Sahel", "en")
    d.cover([
        ("Facility", "OmniIntelOS Niamey DC1"),
        ("Design tier", "Uptime Institute Tier III (concurrently maintainable)"),
        ("Commissioned", "2024-07"),
        ("Design IT load", "1.6 MW"),
        ("Author", "Cloud & Data Centre Operations"),
    ])
    d.add_page()
    d.h1("1. Rationale")
    d.para(
        "Regional customers in banking, public administration and telecommunications increasingly face data "
        "residency obligations that cannot be met by hosting outside the UEMOA zone. DC1 was built to remove "
        "that constraint, and to give OmniIntelOS control over the GPU capacity its computer-vision and NLP "
        "workloads depend on rather than renting it at spot prices abroad.")
    d.h1("2. Design")
    d.kv_table([
        ("Location", "Niamey, Niger - dedicated 33 kV feed"),
        ("Design IT load", "1.6 MW across four halls"),
        ("Redundancy", "N+1 UPS, N+1 CRAC, 2N utility/generator transfer"),
        ("Cooling", "Indirect adiabatic, elevated ASHRAE A3 envelope for the Sahel climate"),
        ("Design PUE", "1.38 annualised (measured 1.41 in first full year)"),
        ("GPU capacity", "224 accelerators across training and inference pools"),
        ("Connectivity", "Three carriers, diverse physical entry, 400 Gbps aggregate"),
        ("Certifications", "ISO 27001 (2024), ISO 14001 (2026), PCI-DSS scoped zone"),
    ])
    d.h1("3. Climate and power reality")
    d.para(
        "Ambient temperatures above 45 C for sustained periods and a grid with historically variable "
        "availability drove two decisions that differ from a temperate-climate design: an elevated thermal "
        "envelope operated to ASHRAE A3 rather than A1, which trades a small amount of hardware lifetime for "
        "a materially lower cooling load, and on-site generation sized for full IT load rather than for a "
        "short bridge to utility restoration.")
    d.para(
        "The Scope 1 emissions in the ESG series step up visibly from 2024 onward for exactly this reason - "
        "generator runtime is a real, measured cost of operating in this location, and the sustainability "
        "roadmap addresses it with solar plus storage rather than by excluding it from the boundary.")
    d.h1("4. Measured operating data")
    ms = [m for m in months() if m >= "2024-07"]
    d.table(["Period", "Energy (kWh)", "Renewable %", "Scope 2 (tCO2e)", "Uptime %", "P99 (ms)"],
            [[m, f"{by[m]['Energy Consumption']:,.0f}", f"{by[m]['Renewable Energy Ratio']:.1f}",
              f"{by[m]['Scope 2 Emissions']:.1f}", f"{by[m]['System Uptime']:.3f}",
              f"{by[m]['API P99 Latency']:.0f}"] for m in ms[::2]],
            [26, 34, 28, 34, 30, 24])
    d.h1("5. Roadmap")
    d.bullets([
        "DC2 Dakar - second availability zone, target commissioning 2027, enabling active/active regional failover.",
        "Solar plus battery storage phase 2, lifting renewable share past 70% and reducing generator runtime.",
        "Liquid cooling pilot for the training pool as accelerator density rises.",
        "Carbon-aware scheduling: shifting non-urgent training to high-renewable windows.",
    ])
    p = out / "omniintelos_dc1_whitepaper_en.pdf"
    d.output(str(p))
    return p


def board_pack(by, year: int, q: int, out: Path) -> Path:
    qm = {1: ["01", "02", "03"], 2: ["04", "05", "06"], 3: ["07", "08", "09"], 4: ["10", "11", "12"]}[q]
    ms = [f"{year}-{x}" for x in qm if f"{year}-{x}" in by]
    if not ms:
        return None
    rev = _agg(by, ms, "Revenue")
    ebitda = _agg(by, ms, "EBITDA")
    hs, hl = health_index(ms[-1], by)
    ph = phase_for(ms[-1])
    d = Doc(f"Board Pack Q{q} {year}", f"Quarter ended {ms[-1]}", "en")
    d.cover([
        ("Reporting period", f"Q{q} {year} ({ms[0]} to {ms[-1]})"),
        ("Prepared for", "Board of Directors, OmniIntelOS S.A."),
        ("Prepared by", "Office of the CFO"),
        ("Enterprise health", f"{hs:.1f}/100 ({hl})"),
        ("Operating regime", ph["label_en"]),
    ])
    d.add_page()
    d.h1("Position")
    d.para(ph["narrative_en"])
    d.kv_table([
        ("Quarterly revenue", usd(rev)),
        ("EBITDA", f"{usd(ebitda)} ({ebitda/rev*100 if rev else 0:.1f}%)"),
        ("Exit ARR", usd(_agg(by, ms, "ARR", "last"))),
        ("NRR (avg)", f"{_agg(by, ms, 'Net Revenue Retention', 'avg'):.1f}%"),
        ("Rule of 40", f"{_agg(by, ms, 'Rule of 40', 'avg'):.1f}%"),
        ("Cash balance", usd(_agg(by, ms, "Cash Balance", "last"))),
        ("Cash runway", f"{_agg(by, ms, 'Cash Runway', 'last'):.1f} months"),
        ("Headcount", f"{_agg(by, ms, 'Headcount', 'last'):,.0f}"),
        ("Uptime (avg)", f"{_agg(by, ms, 'System Uptime', 'avg'):.3f}%"),
        ("Enterprise health", f"{hs:.1f} ({hl})"),
    ])
    d.h1("Monthly detail")
    d.table(["Month", "Revenue", "GM %", "EBITDA %", "ARR", "Churn %", "Uptime %", "Health"],
            [[m, f"{by[m]['Revenue']:,.0f}", f"{by[m]['Gross Margin']:.1f}",
              f"{by[m]['EBITDA Margin']:.1f}", f"{by[m]['ARR']:,.0f}",
              f"{by[m]['Monthly Churn Rate']:.2f}", f"{by[m]['System Uptime']:.3f}",
              f"{health_index(m, by)[0]:.0f}"] for m in ms],
            [22, 30, 20, 24, 30, 22, 24, 20])
    d.h1("Domain analysis")
    d.para(
        "Each line is computed from the quarter's monthly series and judged against the thresholds "
        "published in the domain specification. Statements of the form 'below target' are therefore "
        "checkable against both the stored KPI value and the published band.")
    for dom, mets in DOMAIN_METRICS.items():
        d.h2(dom)
        d.table(["Metric", "Quarter close", "Standing", "Published band"],
                [[m, fmt_metric(m, by[ms[-1]][m]), verdict(m, by[ms[-1]][m]), target_text(m)]
                 for m in mets if m in by[ms[-1]]],
                [62, 32, 34, 58])
        for m in mets:
            c = commentary(by, ms, m)
            if c:
                d.para(c)

    d.h1("Segment revenue")
    d.table(["Region", "Revenue (USD)", "%"],
            [[rg, f"{rev*sh:,.0f}", f"{sh*100:.1f}%"] for rg, sh in REGION_MIX.items()],
            [66, 60, 40])
    d.table(["Service line", "Revenue (USD)", "%"],
            [[ln, f"{rev*sh:,.0f}", f"{sh*100:.1f}%"] for ln, sh in LINE_MIX.items()],
            [86, 50, 38])

    d.h1("Matters for the board")
    if hl in ("At Risk", "Critical"):
        d.callout("Escalation",
                  f"Enterprise health closed the quarter at {hs:.1f} ({hl}). The operating regime is "
                  f"'{ph['label_en']}'. Management requests board direction on the mitigations set out below.",
                  (254, 226, 226))
    d.bullets([
        f"Growth: NRR at {_agg(by, ms, 'Net Revenue Retention', 'avg'):.1f}% against a 110% target.",
        f"Margin: gross margin at {_agg(by, ms, 'Gross Margin', 'avg'):.1f}%, EBITDA margin at {ebitda/rev*100 if rev else 0:.1f}%.",
        f"Liquidity: {_agg(by, ms, 'Cash Runway', 'last'):.1f} months of runway at the current burn.",
        f"People: annualised turnover {_agg(by, ms, 'Annual Employee Turnover', 'avg'):.1f}%, eNPS {_agg(by, ms, 'Employee Net Promoter Score', 'avg'):+.0f}.",
        f"Security: {_agg(by, ms, 'Critical Vulnerabilities', 'max'):.0f} critical vulnerabilities at peak; SLA compliance {_agg(by, ms, 'SLA Compliance', 'avg'):.1f}%.",
    ])
    p = out / f"omniintelos_board_pack_{year}Q{q}_en.pdf"
    d.output(str(p))
    return p


def esg_report(by, year: int, out: Path) -> Path:
    ms = _fy_months(year)
    d = Doc(f"Sustainability & ESG Report {year}", "GHG Protocol Scope 1, 2 and 3 disclosure", "en")
    d.cover([
        ("Reporting year", str(year)),
        ("Boundary", "Operational control, all group entities"),
        ("Standard", "GHG Protocol Corporate Standard"),
        ("Assurance", "Limited assurance, independent third party"),
        ("Author", "Security & Compliance / Facilities"),
    ])
    d.add_page()
    d.h1("1. Basis of preparation")
    d.para(
        "Emissions are reported under the GHG Protocol Corporate Standard using an operational-control "
        "boundary. Scope 1 covers on-site generator fuel and the company vehicle fleet. Scope 2 is "
        "location-based, using a 0.62 tCO2e/MWh grid factor for the Sahel interconnection. Scope 3 covers "
        "purchased goods and services, business travel and upstream cloud capacity.")
    d.para(
        "Generator runtime is inside the boundary and is the single largest driver of the Scope 1 increase "
        "from 2024. Excluding it would make the numbers look better and mean less.")
    d.h1("2. Emissions")
    d.table(["Month", "Scope 1", "Scope 2", "Scope 3", "Total tCO2e", "Energy kWh", "Renewable %"],
            [[m, f"{by[m]['Scope 1 Emissions']:.1f}", f"{by[m]['Scope 2 Emissions']:.1f}",
              f"{by[m]['Scope 3 Emissions']:.1f}", f"{by[m]['Total Carbon Footprint']:.1f}",
              f"{by[m]['Energy Consumption']:,.0f}", f"{by[m]['Renewable Energy Ratio']:.1f}"] for m in ms],
            [22, 24, 24, 24, 28, 30, 26])
    d.h1("3. Intensity")
    d.kv_table([
        ("Total emissions", f"{_agg(by, ms, 'Total Carbon Footprint'):,.0f} tCO2e"),
        ("Carbon intensity", f"{_agg(by, ms, 'Carbon Intensity per Revenue', 'avg'):,.1f} tCO2e per USD million"),
        ("Energy consumed", f"{_agg(by, ms, 'Energy Consumption'):,.0f} kWh"),
        ("Renewable share (exit)", f"{_agg(by, ms, 'Renewable Energy Ratio', 'last'):.1f}%"),
        ("Water consumed", f"{_agg(by, ms, 'Water Consumption'):,.0f} m3"),
        ("Waste diverted (exit)", f"{_agg(by, ms, 'Waste Diverted from Landfill', 'last'):.1f}%"),
    ])
    d.h1("4. Social and governance")
    d.kv_table([
        ("Headcount (exit)", f"{_agg(by, ms, 'Headcount', 'last'):,.0f}"),
        ("Board diversity ratio", f"{_agg(by, ms, 'Board Diversity Ratio', 'last'):.1f}%"),
        ("Audit compliance score", f"{_agg(by, ms, 'Audit Compliance Score', 'last'):.1f}%"),
        ("Privacy incidents", f"{_agg(by, ms, 'Privacy Incident Count'):.0f}"),
        ("Employee turnover", f"{_agg(by, ms, 'Annual Employee Turnover', 'avg'):.1f}%"),
        ("Training hours per employee", f"{_agg(by, ms, 'Training Hours per Employee', 'avg'):.1f}"),
    ])
    d.h1("5. Methodology and emission factors")
    d.kv_table([
        ("Standard", "GHG Protocol Corporate Accounting and Reporting Standard"),
        ("Consolidation approach", "Operational control"),
        ("Grid emission factor", "0.62 tCO2e/MWh (Sahel interconnection, location-based)"),
        ("Generator fuel factor", "2.68 kgCO2e per litre diesel"),
        ("Scope 3 categories included", "1 (purchased goods), 6 (business travel), 8 (upstream cloud)"),
        ("Scope 3 categories excluded", "11 (use of sold products) - not yet measurable with confidence"),
        ("Base year", "2020"),
        ("Assurance", "Limited assurance, independent third party"),
        ("Restatement policy", "Restate if a boundary change moves base-year emissions by over 5%"),
    ])
    d.para(
        "Where a factor is uncertain, the more conservative (higher-emission) value is used. Categories "
        "that cannot yet be measured with confidence are named and excluded rather than estimated to a "
        "number that would look complete and mean nothing.")

    d.h1("6. Targets and roadmap")
    d.table(["Target", "Baseline", "Current", "Target date"], [
        ["Renewable energy share >= 60%", "12% (2020)", f"{_agg(by, ms, 'Renewable Energy Ratio', 'last'):.1f}%", "2027"],
        ["Board diversity >= 40%", "14% (2020)", f"{_agg(by, ms, 'Board Diversity Ratio', 'last'):.1f}%", "2027"],
        ["Audit compliance >= 98%", "86% (2020)", f"{_agg(by, ms, 'Audit Compliance Score', 'last'):.1f}%", "2026"],
        ["Zero privacy incidents", "3 (2023)", f"{_agg(by, ms, 'Privacy Incident Count'):.0f}", "Ongoing"],
        ["Waste diverted >= 75%", "31% (2020)", f"{_agg(by, ms, 'Waste Diverted from Landfill', 'last'):.1f}%", "2028"],
        ["ISO 14001 certification", "Not held", "In progress", "2026"],
    ], [70, 34, 34, 34])

    d.h1("7. Commentary")
    for met in DOMAIN_METRICS["ESG"]:
        c = commentary(by, ms, met)
        if c:
            d.para(c)

    d.h1("8. Appendix - monthly ESG series")
    esg_mets = ["Scope 1 Emissions", "Scope 2 Emissions", "Scope 3 Emissions",
                "Total Carbon Footprint", "Energy Consumption", "Renewable Energy Ratio",
                "Water Consumption", "Waste Diverted from Landfill", "Board Diversity Ratio",
                "Audit Compliance Score", "Privacy Incident Count", "Carbon Intensity per Revenue"]
    for chunk_start in range(0, len(esg_mets), 4):
        chunk = esg_mets[chunk_start:chunk_start + 4]
        d.table(["Period"] + [m[:26] for m in chunk],
                [[m] + [f"{by[m].get(x, 0):,.2f}" for x in chunk] for m in ms],
                [24] + [41] * len(chunk))

    p = out / f"omniintelos_esg_report_{year}_en.pdf"
    d.output(str(p))
    return p


# ─────────────────────────────────────────────────────────────────────────────
# Meeting minutes / notes / policies (markdown - the RAG's best substrate)
# ─────────────────────────────────────────────────────────────────────────────

def meeting_minutes(by, out: Path) -> List[Path]:
    written = []
    picks = [
        ("2020-05", "fr", "Comite de crise COVID-19"),
        ("2021-08", "en", "Series A closing and use of proceeds"),
        ("2022-09", "en", "Talent retention emergency review"),
        ("2023-02", "en", "SEV-1 incident bridge - INC-2023-0214"),
        ("2023-03", "fr", "Comite de direction - suites de l'incident"),
        ("2024-03", "fr", "Revue d'avancement DC1 Niamey"),
        ("2024-10", "en", "Rule of 40 efficiency programme kickoff"),
        ("2025-05", "fr", "Cellule logistique - ruptures d'approvisionnement Sahel"),
        ("2025-11", "en", "Generative AI demand review"),
        ("2026-04", "en", "ISO 14001 readiness and ESG targets"),
    ]
    for period, lang, subject in picks:
        if period not in by:
            continue
        k = by[period]
        ph = phase_for(period)
        hs, hl = health_index(period, by)
        fr = lang == "fr"
        att = ("Presents" if fr else "Attendees")
        lines = [
            f"# {subject}",
            "",
            f"**{'Entreprise' if fr else 'Company'}:** {COMPANY} - {HQ}  ",
            f"**{'Periode' if fr else 'Period'}:** {period}  ",
            f"**{'Regime operationnel' if fr else 'Operating regime'}:** {ph['label_fr'] if fr else ph['label_en']}  ",
            f"**{'Indice de sante' if fr else 'Enterprise health'}:** {hs:.1f}/100 ({hl})  ",
            f"**{att}:** Direction generale, Finance, Ingenierie, Securite, People & Culture"
            if fr else f"**{att}:** CEO, CFO, CTO, CISO, Chief People Officer",
            "",
            "## " + ("Contexte" if fr else "Context"),
            "",
            ph["narrative_fr"] if fr else ph["narrative_en"],
            "",
            "## " + ("Indicateurs examines" if fr else "Metrics reviewed"),
            "",
            f"| {'Indicateur' if fr else 'Metric'} | {'Valeur' if fr else 'Value'} |",
            "|---|---|",
            f"| {'Chiffre d affaires mensuel' if fr else 'Monthly revenue'} | {usd(k['Revenue'])} |",
            f"| ARR | {usd(k['ARR'])} |",
            f"| {'Marge brute' if fr else 'Gross margin'} | {k['Gross Margin']:.1f}% |",
            f"| NRR | {k['Net Revenue Retention']:.1f}% |",
            f"| {'Attrition client' if fr else 'Logo churn'} | {k['Monthly Churn Rate']:.2f}% |",
            f"| {'Disponibilite' if fr else 'Uptime'} | {k['System Uptime']:.3f}% |",
            f"| {'Vulnerabilites critiques' if fr else 'Critical vulnerabilities'} | {k['Critical Vulnerabilities']:.0f} |",
            f"| OEE | {k['Overall Equipment Effectiveness']:.1f}% |",
            f"| {'Livraison a temps' if fr else 'On-time delivery'} | {k['On-Time Delivery Rate']:.1f}% |",
            f"| {'Effectif' if fr else 'Headcount'} | {k['Headcount']:,.0f} |",
            f"| eNPS | {k['Employee Net Promoter Score']:+.0f} |",
            f"| {'Autonomie de tresorerie' if fr else 'Cash runway'} | {k['Cash Runway']:.1f} {'mois' if fr else 'months'} |",
            "",
            "## " + ("Analyse par domaine" if fr else "Domain analysis"),
            "",
            ("Chaque ligne est calculee a partir des donnees du mois et comparee aux seuils publies "
             "dans la specification des domaines." if fr else
             "Each line is computed from the month's data and compared against the thresholds published "
             "in the domain specification."),
            "",
        ]
        for dom, mets in DOMAIN_METRICS.items():
            lines.append(f"### {dom}")
            lines.append("")
            for met in mets:
                if met in k:
                    lines.append(f"- {commentary(by, [period], met, lang)}")
            lines.append("")
        lines += [
            "## " + ("Decisions" if fr else "Decisions"),
            "",
        ]
        decisions_fr = {
            "2020-05": ["Gel des recrutements non critiques jusqu'au retablissement du pipeline.",
                        "Passage integral en livraison a distance pour les projets en cours.",
                        "Negociation de delais de paiement avec les fournisseurs materiels."],
            "2023-03": ["Financement immediat du programme de durcissement zero-trust.",
                        "Communication proactive aux clients affectes, en francais et en anglais.",
                        "Lancement de la certification ISO 27001 sous pilotage du comite de direction."],
            "2024-03": ["Validation du calendrier de mise en service du DC1 pour juillet 2024.",
                        "Approbation de la facilite d'equipement pour le materiel importe.",
                        "Plan de secours energetique valide avec dimensionnement pleine charge."],
            "2025-05": ["Constitution de stocks tampons pour les composants critiques.",
                        "Diversification des corridors logistiques via Cotonou et Lome.",
                        "Revision des engagements de delai client tant que la perturbation dure."],
        }
        decisions_en = {
            "2021-08": ["Allocate USD 12M Series A: 45% engineering hiring, 30% DC1 land and design, 25% working capital.",
                        "Open the Dakar commercial office in Q4 2021.",
                        "Institute quarterly board reporting on the Rule of 40."],
            "2022-09": ["Immediate market salary review for engineering and data science bands.",
                        "Cap onboarding intake at a level the mentoring pool can absorb.",
                        "Introduce structured career frameworks to address the eNPS decline."],
            "2023-02": ["Isolate the staging estate from the transit VPC immediately, accepting customer-facing degradation.",
                        "Retain external DFIR; begin forensic imaging.",
                        "Enforce MFA across every environment without waiting for the post-mortem.",
                        "Prepare regulator notification inside the 72-hour window."],
            "2024-10": ["Target Rule of 40 above 40% within two quarters.",
                        "Automate delivery tooling to lift gross margin without headcount reduction.",
                        "Revise pricing on legacy contracts at renewal."],
            "2025-11": ["Prioritise bilingual FR/EN document intelligence in the roadmap.",
                        "Reserve GPU capacity in DC1 for inference demand ahead of training.",
                        "Expand the partner channel across academia and regional banks."],
            "2026-04": ["Commit to ISO 14001 certification within the calendar year.",
                        "Approve solar plus storage phase 2 for DC1.",
                        "Set a board diversity target of 45% by end 2027."],
        }
        for dline in (decisions_fr if fr else decisions_en).get(period, [
            "Poursuite du plan en cours, revue au prochain comite." if fr else
            "Continue the current plan; review at the next committee."]):
            lines.append(f"- {dline}")
        lines += ["", "## " + ("Actions" if fr else "Actions"), "",
                  f"| {'Action' if fr else 'Action'} | {'Responsable' if fr else 'Owner'} | {'Echeance' if fr else 'Due'} |",
                  "|---|---|---|",
                  f"| {'Suivi des indicateurs ci-dessus' if fr else 'Track the metrics above'} | Finance | {period}-28 |",
                  f"| {'Rapport au conseil' if fr else 'Report to the board'} | CEO | {period}-30 |",
                  "",
                  "---", "",
                  ("*OmniIntelOS S.A. est une entreprise fictive ; ce compte rendu est genere pour la "
                   "demonstration d'IntelAI.*" if fr else
                   "*OmniIntelOS S.A. is a fictional company; these minutes are generated for IntelAI "
                   "demonstration purposes.*")]
        p = out / f"omniintelos_minutes_{period}_{lang}.md"
        p.write_text("\n".join(lines), encoding="utf-8")
        written.append(p)
    return written


def policies(out: Path) -> List[Path]:
    docs = {
        "omniintelos_policy_information_security_en.md": f"""# Information Security Policy

**{COMPANY}** | Owner: Security & Compliance | Version 4.2 | Effective 2026-01-01

## 1. Purpose and scope
This policy applies to all employees, contractors and third parties with access to
{COMPANY_SHORT} systems, across every entity and every environment.

## 2. Access control
- Multi-factor authentication is mandatory on **all** environments. There is no
  non-production exemption. This requirement originates from incident INC-2023-0214,
  in which a staging estate exempted from MFA was the initial access vector.
- Access is granted on least privilege and reviewed quarterly.
- Joiner-mover-leaver applies identically to employees, contractors and service accounts.
  Access is revoked no later than the last working day.

## 3. Data classification and handling
| Class | Examples | Storage rule |
|---|---|---|
| Restricted | Customer telemetry, credentials, personal data | Production only, encrypted at rest and in transit |
| Confidential | Financial results pre-announcement, contracts | Group systems only, need-to-know |
| Internal | Runbooks, minutes, architecture | Group systems, all staff |
| Public | Marketing, published research | No restriction |

**Customer data must never be replicated into a non-production environment.** An automated
scanner verifies this daily and raises a SEV-2 on detection.

## 4. Network and estate segmentation
Production, staging, development and corporate estates are segmented under a zero-trust model.
No implicit trust is granted on the basis of network location.

## 5. Vulnerability management
| Severity (CVSS v3.1) | Remediation SLA |
|---|---|
| Critical (9.0-10.0) | 72 hours |
| High (7.0-8.9) | 14 days |
| Medium (4.0-6.9) | 60 days |
| Low (0.1-3.9) | Next maintenance window |

## 6. Incident response
Severity levels, response targets and on-call compensation are defined in the Employee
Handbook Annex B. All SEV-1 incidents require a written post-mortem within 10 working days,
published internally without redaction of root cause.

## 7. Exceptions
Exceptions require written CISO approval, a compensating control and an expiry date. No
exception may exceed 90 days without board-level review.
""",
        "omniintelos_politique_protection_donnees_fr.md": f"""# Politique de protection des donnees personnelles

**{COMPANY}** | Responsable : Securite et Conformite | Version 3.1 | En vigueur au 2026-01-01

## 1. Objet
La presente politique definit le traitement des donnees a caractere personnel par
{COMPANY_SHORT} dans l'ensemble de ses entites, en Afrique de l'Ouest, en Europe et en
Amerique du Nord.

## 2. Cadre applicable
- Loi nigerienne relative a la protection des donnees a caractere personnel
- Acte additionnel A/SA.1/01/10 de la CEDEAO
- Reglement general sur la protection des donnees (RGPD) pour les traitements europeens
- Exigences contractuelles specifiques des clients du secteur bancaire

## 3. Principes
- **Minimisation** : seules les donnees necessaires a la finalite declaree sont collectees.
- **Limitation de conservation** : la telemetrie client est conservee 24 mois, les journaux
  de securite 12 mois, les donnees de facturation 10 ans (obligation comptable OHADA).
- **Localisation** : les clients soumis a une obligation de residence des donnees sont
  heberges exclusivement au DC1 de Niamey.

## 4. Droits des personnes concernees
Les demandes d'acces, de rectification, d'effacement et de portabilite sont traitees dans un
delai de 30 jours. Un registre des demandes est tenu et audite trimestriellement.

## 5. Violations de donnees
Toute violation est notifiee a l'autorite de controle competente dans les 72 heures et aux
personnes concernees sans delai injustifie lorsqu'un risque eleve existe. L'incident
INC-2023-0214 a donne lieu a trois notifications formelles.

## 6. Sous-traitants
Tout sous-traitant fait l'objet d'une evaluation prealable et d'un accord de traitement
comportant des clauses de securite et d'audit.
""",
        "omniintelos_runbook_sev1_en.md": f"""# Runbook: SEV-1 incident response

**{COMPANY}** | Cloud & Data Centre Operations | Reviewed 2026-02

## Declaration
Any of the following declares a SEV-1 immediately, without waiting for confirmation:
- Customer-facing unavailability affecting more than one tenant
- Confirmed or suspected unauthorised access to any estate
- Data loss or suspected exfiltration
- Loss of both utility and generator power at DC1

## First 15 minutes
1. Page the on-call SRE and the incident commander rota.
2. Open the incident channel and the incident document. One writer, everyone else reads.
3. Assign three roles explicitly: Incident Commander, Communications Lead, Scribe.
4. State the current hypothesis and the next check that would disprove it.

## Containment principles
- Containment beats diagnosis. Isolate first, understand afterwards.
- Accept customer-facing degradation to stop an active intrusion. This is a standing
  decision made by the board following INC-2023-0214 and does not require re-approval.
- Preserve forensic state before remediation: image before you rebuild.

## Communications
| Audience | Trigger | Owner | Channel |
|---|---|---|---|
| Affected customers | Confirmed impact | Communications Lead | Status page + direct account contact, EN and FR |
| Regulator | Confirmed personal data breach | DPO | Formal notification within 72h |
| Board | Any SEV-1 lasting over 4h | CEO | Direct briefing |

## Closure
An incident is not closed when service is restored. It is closed when the written
post-mortem is published and every remediation action has an owner and a date.
""",
        "omniintelos_architecture_decision_records_en.md": f"""# Architecture Decision Records (selected)

**{COMPANY}** | Engineering | Maintained continuously

## ADR-011: Bilingual-first document intelligence
**Status:** Accepted (2021-04) | **Context:** Regional customers operate in French while
international partners operate in English, and documents routinely mix both.
**Decision:** Every retrieval and extraction component must be evaluated on a bilingual
FR/EN corpus before release. A model that scores well only in English is not shippable.
**Consequence:** Slower model iteration; a durable commercial moat in francophone markets.

## ADR-017: Own the inference capacity
**Status:** Accepted (2023-09) | **Context:** GPU spot pricing abroad is volatile and
several customers face data-residency obligations that offshore hosting cannot meet.
**Decision:** Build and operate DC1 in Niamey rather than renting capacity.
**Consequence:** Heavy CAPEX and leverage through 2024; sovereign-hosting revenue and
predictable inference cost from 2024-07.

## ADR-021: Zero trust between estates
**Status:** Accepted (2023-03) | **Context:** INC-2023-0214 showed that network location
had been treated as an implicit trust boundary.
**Decision:** No estate trusts another by virtue of network position. All inter-estate
traffic is authenticated and authorised per request.
**Consequence:** Higher operational complexity; the class of failure that produced
INC-2023-0214 is structurally removed.

## ADR-026: Carbon-aware batch scheduling
**Status:** Proposed (2026-03) | **Context:** Scope 2 emissions track grid mix, which varies
predictably with solar availability at DC1.
**Decision:** Defer non-urgent training workloads into high-renewable windows.
**Consequence:** Longer worst-case training latency; measurable Scope 2 reduction.
""",
    }
    written = []
    for name, body in docs.items():
        p = out / name
        p.write_text(body, encoding="utf-8")
        written.append(p)
    return written


# ─────────────────────────────────────────────────────────────────────────────
# Spreadsheets, decks, charts
# ─────────────────────────────────────────────────────────────────────────────

def workbooks(by, out: Path) -> List[Path]:
    if not _XLSX:
        return []
    ms = months()
    written = []

    wb = openpyxl.Workbook()
    hdr_font = Font(bold=True, color="FFFFFF")
    hdr_fill = PatternFill("solid", fgColor="2563EB")

    domains = ["Finance", "Growth", "People", "Operations", "Logistics", "IT", "ESG"]
    kpis = generate_kpis()
    per_domain: Dict[str, List[str]] = {}
    for r in kpis:
        per_domain.setdefault(r["category"], [])
        if r["metric_name"] not in per_domain[r["category"]]:
            per_domain[r["category"]].append(r["metric_name"])

    first = True
    for dom in domains:
        ws = wb.active if first else wb.create_sheet()
        ws.title = dom
        first = False
        mets = per_domain.get(dom, [])
        ws.cell(1, 1, "Period").font = hdr_font
        ws.cell(1, 1).fill = hdr_fill
        for j, met in enumerate(mets, start=2):
            c = ws.cell(1, j, met)
            c.font = hdr_font
            c.fill = hdr_fill
            ws.column_dimensions[c.column_letter].width = max(14, min(30, len(met) + 3))
        ws.column_dimensions["A"].width = 11
        for i, m in enumerate(ms, start=2):
            ws.cell(i, 1, m)
            for j, met in enumerate(mets, start=2):
                ws.cell(i, j, round(by[m].get(met, 0.0), 3))
        ws.freeze_panes = "B2"
    p = out / "omniintelos_kpi_workbook_2020_2026.xlsx"
    wb.save(p)
    written.append(p)

    wb2 = openpyxl.Workbook()
    ws = wb2.active
    ws.title = "P&L monthly"
    cols = ["Period", "Revenue", "Subscription", "Services", "Data Centre", "COGS", "Gross Profit",
            "Gross Margin %", "Opex", "EBITDA", "EBITDA %", "D&A", "Interest", "Taxes",
            "Net Income", "Net Margin %", "Cash", "Debt", "Equity", "D/E", "Runway (mo)"]
    for j, c in enumerate(cols, 1):
        cc = ws.cell(1, j, c)
        cc.font = hdr_font
        cc.fill = hdr_fill
        ws.column_dimensions[cc.column_letter].width = 15
    keys = ["Revenue", "Subscription Revenue", "Professional Services Revenue", "Data Centre Revenue",
            "COGS", "Gross Profit", "Gross Margin", "Operating Expenses", "EBITDA", "EBITDA Margin",
            "Depreciation & Amortisation", "Interest Expense", "Taxes", "Net Income",
            "Net Profit Margin", "Cash Balance", "Total Debt", "Shareholders Equity",
            "Debt to Equity", "Cash Runway"]
    for i, m in enumerate(ms, start=2):
        ws.cell(i, 1, m)
        for j, k in enumerate(keys, start=2):
            ws.cell(i, j, round(by[m].get(k, 0.0), 2))
    ws.freeze_panes = "B2"
    ws2 = wb2.create_sheet("Assumptions")
    for i, (k, v) in enumerate([
        ("Company", COMPANY), ("Headquarters", HQ), ("Status", "FICTIONAL - generated dataset"),
        ("Functional currency", "XOF"), ("Presentation currency", "USD"),
        ("XOF per EUR (BCEAO fixed peg)", 655.957), ("XOF per USD (planning)", round(XOF_PER_USD, 3)),
        ("Corporate income tax", "30% (Niger)"), ("Grid emission factor", "0.62 tCO2e/MWh"),
        ("Period covered", f"{ms[0]} to {ms[-1]} ({len(ms)} months)"),
        ("Gross margin formula", "(Revenue - COGS) / Revenue"),
        ("Rule of 40 formula", "YoY revenue growth % + EBITDA margin %"),
    ], start=1):
        ws2.cell(i, 1, k).font = Font(bold=True)
        ws2.cell(i, 2, v)
    ws2.column_dimensions["A"].width = 34
    ws2.column_dimensions["B"].width = 46
    p2 = out / "omniintelos_financial_model_2020_2026.xlsx"
    wb2.save(p2)
    written.append(p2)

    wb3 = openpyxl.Workbook()
    ws = wb3.active
    ws.title = "Headcount plan"
    for j, c in enumerate(["Period", "Headcount", "New Hires", "Separations", "Turnover %",
                           "Open Positions", "Time to Hire", "eNPS", "Revenue per Employee",
                           "Cost Per Hire"], 1):
        cc = ws.cell(1, j, c)
        cc.font = hdr_font
        cc.fill = hdr_fill
        ws.column_dimensions[cc.column_letter].width = 18
    hk = ["Headcount", "New Hires", "Separations", "Annual Employee Turnover", "Open Positions",
          "Time to Hire", "Employee Net Promoter Score", "Revenue per Employee", "Cost Per Hire"]
    for i, m in enumerate(ms, start=2):
        ws.cell(i, 1, m)
        for j, k in enumerate(hk, start=2):
            ws.cell(i, j, round(by[m].get(k, 0.0), 2))
    ws2 = wb3.create_sheet("Departments")
    ws2.cell(1, 1, "Department").font = Font(bold=True)
    ws2.cell(1, 2, "Share of headcount").font = Font(bold=True)
    shares = [0.27, 0.16, 0.13, 0.14, 0.09, 0.08, 0.05, 0.04, 0.04]
    for i, (dep, sh) in enumerate(zip(DEPARTMENTS, shares), start=2):
        ws2.cell(i, 1, dep)
        ws2.cell(i, 2, sh)
    ws2.column_dimensions["A"].width = 34
    ws2.column_dimensions["B"].width = 20
    p3 = out / "omniintelos_headcount_plan_2020_2026.xlsx"
    wb3.save(p3)
    written.append(p3)
    return written


def charts(by, out: Path) -> List[Path]:
    if not _MPL:
        return []
    ms = months()
    written = []
    specs = [
        ("arr_growth", "ARR (USD)", [by[m]["ARR"] for m in ms], "#2563eb"),
        ("gross_margin", "Gross margin (%)", [by[m]["Gross Margin"] for m in ms], "#059669"),
        ("uptime", "System uptime (%)", [by[m]["System Uptime"] for m in ms], "#dc2626"),
        ("headcount", "Headcount", [by[m]["Headcount"] for m in ms], "#7c3aed"),
        ("health", "Enterprise health index", [health_index(m, by)[0] for m in ms], "#ea580c"),
        ("emissions", "Total carbon footprint (tCO2e)", [by[m]["Total Carbon Footprint"] for m in ms], "#0891b2"),
    ]
    for name, label, series, color in specs:
        fig, ax = plt.subplots(figsize=(9, 3.6), dpi=150)
        ax.plot(range(len(ms)), series, color=color, linewidth=1.9)
        ax.fill_between(range(len(ms)), series, alpha=0.12, color=color)
        ax.set_title(f"{COMPANY_SHORT} - {label}  ({ms[0]} to {ms[-1]})", fontsize=11, loc="left")
        ticks = list(range(0, len(ms), 6))
        ax.set_xticks(ticks)
        ax.set_xticklabels([ms[i] for i in ticks], rotation=45, ha="right", fontsize=7)
        ax.grid(alpha=0.25, linewidth=0.6)
        ax.spines[["top", "right"]].set_visible(False)
        # Mark the breach so the chart is self-explaining
        if "2023-02" in ms:
            bi = ms.index("2023-02")
            ax.axvline(bi, color="#991b1b", linestyle="--", linewidth=1, alpha=0.7)
            ax.annotate("INC-2023-0214", xy=(bi, max(series)), fontsize=6.5,
                        color="#991b1b", rotation=90, va="top", ha="right")
        fig.tight_layout()
        p = out / f"omniintelos_chart_{name}.png"
        fig.savefig(p)
        plt.close(fig)
        written.append(p)
    return written


def decks(by, out: Path) -> List[Path]:
    if not _PPTX:
        return []
    written = []
    for year, q in [(2023, 1), (2025, 4)]:
        qm = {1: ["01", "02", "03"], 4: ["10", "11", "12"]}[q]
        ms = [f"{year}-{x}" for x in qm if f"{year}-{x}" in by]
        if not ms:
            continue
        hs, hl = health_index(ms[-1], by)
        ph = phase_for(ms[-1])
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[0])
        s.shapes.title.text = f"{COMPANY_SHORT} - Board Review Q{q} {year}"
        s.placeholders[1].text = f"{ph['label_en']}  |  Enterprise health {hs:.1f}/100 ({hl})"

        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Position"
        s.placeholders[1].text_frame.text = ph["narrative_en"]

        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Key metrics"
        tf = s.placeholders[1].text_frame
        rows = [
            f"Revenue (quarter): {usd(_agg(by, ms, 'Revenue'))}",
            f"Exit ARR: {usd(_agg(by, ms, 'ARR', 'last'))}",
            f"Gross margin: {_agg(by, ms, 'Gross Margin', 'avg'):.1f}%",
            f"NRR: {_agg(by, ms, 'Net Revenue Retention', 'avg'):.1f}%",
            f"Rule of 40: {_agg(by, ms, 'Rule of 40', 'avg'):.1f}%",
            f"Uptime: {_agg(by, ms, 'System Uptime', 'avg'):.3f}%",
            f"Headcount: {_agg(by, ms, 'Headcount', 'last'):,.0f}",
            f"Cash runway: {_agg(by, ms, 'Cash Runway', 'last'):.1f} months",
        ]
        tf.text = rows[0]
        for r in rows[1:]:
            tf.add_paragraph().text = r

        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = "Matters for the board"
        tf = s.placeholders[1].text_frame
        items = ([f"Enterprise health at {hs:.1f} ({hl}) - escalation required",
                  "Security remediation programme funding",
                  "Customer retention response"] if hl in ("At Risk", "Critical") else
                 ["Growth and margin both on target",
                  "Capacity planning for continued demand",
                  "ESG certification roadmap"])
        tf.text = items[0]
        for it in items[1:]:
            tf.add_paragraph().text = it

        p = out / f"omniintelos_board_deck_{year}Q{q}_en.pptx"
        prs.save(p)
        written.append(p)
    return written


# ─────────────────────────────────────────────────────────────────────────────

def build_corpus(root: Path) -> Dict[str, Any]:
    """Generate the whole estate under `root`, returning a manifest."""
    kpis = generate_kpis()
    by = _by_period(kpis)
    root.mkdir(parents=True, exist_ok=True)
    for d in ("Finance", "Growth", "People", "Operations", "Logistics", "IT", "ESG", "Corporate"):
        (root / d).mkdir(exist_ok=True)

    made: Dict[str, List[str]] = {}
    skipped: List[str] = []

    def rec(bucket: str, paths):
        made.setdefault(bucket, [])
        for p in paths:
            if p:
                made[bucket].append(str(p))

    if _FPDF:
        rec("pdf", [annual_report(by, y, "en", root / "Finance") for y in range(2020, 2026)])
        rec("pdf", [annual_report(by, y, "fr", root / "Finance") for y in (2023, 2025)])
        rec("pdf", [incident_postmortem(by, root / "IT")])
        rec("pdf", [employee_handbook(root / "People", "en"), employee_handbook(root / "People", "fr")])
        rec("pdf", [datacentre_whitepaper(by, root / "IT")])
        rec("pdf", [esg_report(by, y, root / "ESG") for y in (2024, 2025)])
        rec("pdf", [board_pack(by, y, q, root / "Corporate")
                    for y, q in [(2020, 2), (2021, 3), (2022, 4), (2023, 1), (2023, 3),
                                 (2024, 2), (2024, 4), (2025, 2), (2025, 4), (2026, 2)]])
    else:
        skipped.append("PDF (fpdf2 not installed)")

    rec("markdown", meeting_minutes(by, root / "Corporate"))
    rec("markdown", policies(root / "Corporate"))

    if _XLSX:
        rec("xlsx", workbooks(by, root / "Finance"))
    else:
        skipped.append("XLSX (openpyxl not installed)")
    if _MPL:
        rec("png", charts(by, root / "Corporate"))
    else:
        skipped.append("PNG charts (matplotlib not installed)")
    if _PPTX:
        rec("pptx", decks(by, root / "Corporate"))
    else:
        skipped.append("PPTX (python-pptx not installed)")
    if _AUDIO:
        rec("wav", build_audio(root))
    else:
        skipped.append("WAV audio (TTS_ENDPOINT_URL/TTS_ENDPOINT_TOKEN not set)")

    total_bytes = sum(Path(f).stat().st_size for fs in made.values() for f in fs)
    return {"files": made, "skipped": skipped,
            "count": sum(len(v) for v in made.values()), "bytes": total_bytes}


if __name__ == "__main__":
    import sys
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("data/omniintelos")
    man = build_corpus(target)
    for bucket, files in sorted(man["files"].items()):
        print(f"{bucket:9} {len(files):>3} files")
        for f in files:
            print(f"          {Path(f).stat().st_size:>9,}b  {Path(f).name}")
    if man["skipped"]:
        print("\nSKIPPED:")
        for s in man["skipped"]:
            print("  -", s)
    print(f"\ntotal {man['count']} files, {man['bytes']:,} bytes")
