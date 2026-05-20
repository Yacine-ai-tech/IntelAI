"""
Board-ready exports — PDF & PowerPoint reports.
"""
from __future__ import annotations

import io
from typing import List

import pandas as pd
import plotly.io as pio

from src.core.i18n import I18N, t
from src.core.logger import get_logger
from src.utils.charts import create_category_breakdown_chart, create_health_gauge, create_risk_radar_chart

log = get_logger(__name__)

try:
    from fpdf import FPDF
    _FPDF = True
except ImportError:
    _FPDF = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    _PPTX = True
except ImportError:
    _PPTX = False


def _plot_to_image(fig) -> io.BytesIO:
    img_bytes = pio.to_image(fig, format="png", width=900, height=520, scale=2)
    return io.BytesIO(img_bytes)


def _lbl(en: str, fr: str) -> str:
    return fr if I18N.lang() == "fr" else en


# ── PDF ───────────────────────────────────────────────────────────────────

def build_pdf_report(health: dict, risk: dict, highlights: List[str], category_df: pd.DataFrame) -> bytes:
    if not _FPDF:
        raise RuntimeError("fpdf2 not installed — pip install fpdf2")

    pdf = FPDF(orientation="P", unit="mm", format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, _lbl("OmniIntelOS — Executive Briefing", "OmniIntelOS — Rapport exécutif"), ln=True)

    pdf.set_font("Helvetica", "", 12)
    pdf.cell(0, 8, f"{_lbl('Health Index', 'Indice de santé')}: {health.get('score', 0):.0f} ({health.get('label', '')})", ln=True)
    pdf.cell(0, 8, f"{_lbl('Risk Score', 'Score de risque')}: {risk.get('score', 0):.0f} ({risk.get('label', '')})", ln=True)

    pdf.ln(4)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, _lbl("Executive Highlights", "Points clés"), ln=True)
    pdf.set_font("Helvetica", "", 11)
    for line in highlights:
        pdf.multi_cell(170, 6, f"– {line}")

    pdf.ln(4)
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 8, _lbl("Category Performance", "Performance par catégorie"), ln=True)

    chart = create_category_breakdown_chart(category_df)
    chart_image = _plot_to_image(chart)
    pdf.image(chart_image, x=15, w=180)

    return bytes(pdf.output(dest="S"))


# ── PowerPoint ────────────────────────────────────────────────────────────

def build_pptx_report(health: dict, risk: dict, highlights: List[str], category_df: pd.DataFrame) -> bytes:
    if not _PPTX:
        raise RuntimeError("python-pptx not installed — pip install python-pptx")

    prs = Presentation()

    # Title slide
    slide0 = prs.slides.add_slide(prs.slide_layouts[0])
    slide0.shapes.title.text = "OmniIntelOS"
    slide0.placeholders[1].text = _lbl("Executive Briefing", "Rapport exécutif")

    # Summary slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])
    slide1.shapes.title.text = _lbl("Health & Risk Summary", "Résumé santé & risque")
    body = slide1.shapes.placeholders[1].text_frame
    body.text = f"{_lbl('Health Index', 'Indice de santé')}: {health.get('score', 0):.0f} ({health.get('label', '')})"
    body.add_paragraph().text = f"{_lbl('Risk Score', 'Score de risque')}: {risk.get('score', 0):.0f} ({risk.get('label', '')})"

    # Highlights slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = _lbl("Executive Highlights", "Points clés")
    hb = slide2.shapes.placeholders[1].text_frame
    if highlights:
        hb.text = highlights[0]
        for line in highlights[1:]:
            hb.add_paragraph().text = line

    # Chart slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = _lbl("Category Performance", "Performance par catégorie")
    chart_img = _plot_to_image(create_category_breakdown_chart(category_df))
    slide3.shapes.add_picture(chart_img, Inches(1), Inches(1.5), width=Inches(8))

    # Gauge + Radar slide
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = _lbl("Health & Risk Radar", "Santé & Radar de risque")
    slide4.shapes.add_picture(_plot_to_image(create_health_gauge(health)), Inches(0.6), Inches(1.6), width=Inches(4))
    slide4.shapes.add_picture(_plot_to_image(create_risk_radar_chart(risk)), Inches(5.0), Inches(1.6), width=Inches(4))

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()
